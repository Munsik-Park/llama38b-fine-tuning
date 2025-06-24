from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[1]
sys.path.append(str(ROOT))

from scripts.extract_with_unstructured import filter_elements

try:
    from unstructured.partition.auto import partition
except Exception as e:  # pragma: no cover - missing dependency
    partition = None


def create_docx(path: Path) -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("Title", level=1)
    doc.add_paragraph("Docx content")
    doc.save(path)


def create_pdf(path: Path) -> None:
    import fitz  # PyMuPDF

    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "PDF content")
    doc.save(path)


def create_pptx(path: Path) -> None:
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide"
    slide.shapes.add_textbox(0, 0, 100, 100).text = "PPTX content"
    prs.save(path)


def run_partition(file_path: Path):
    assert partition is not None, "unstructured not installed"
    elements = partition(filename=str(file_path))
    blocks = filter_elements(elements)
    assert blocks, "no blocks returned"
    for block in blocks:
        assert block["text"].strip() != ""
        assert block["type"] in {"NarrativeText", "Title"}
        if block.get("metadata"):
            assert isinstance(block["metadata"], dict)
    return blocks


def test_extractors(tmp_path: Path) -> None:
    if partition is None:
        # Skip test if unstructured cannot be imported
        import pytest

        pytest.skip("unstructured not installed")

    docx_file = tmp_path / "sample.docx"
    pdf_file = tmp_path / "sample.pdf"
    pptx_file = tmp_path / "sample.pptx"

    create_docx(docx_file)
    create_pdf(pdf_file)
    create_pptx(pptx_file)

    for f in [docx_file, pdf_file, pptx_file]:
        blocks = run_partition(f)
        assert len(blocks) > 0

