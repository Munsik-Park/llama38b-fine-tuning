from __future__ import annotations

from typing import List

from pptx import Presentation

from .base_extractor import BaseExtractor


class PPTExtractor(BaseExtractor):
    """Extract text from PPT/PPTX files using python-pptx."""

    def extract_text(self) -> str:
        texts: List[str] = []
        try:
            prs = Presentation(self.file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:  # pragma: no cover - simple error handling
            return f"Extraction failed: {e}"
